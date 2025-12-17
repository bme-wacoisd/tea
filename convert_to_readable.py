#!/usr/bin/env python3
"""
Convert DOCX and PPTX files to Markdown format for easy reading.
Supports batch conversion of entire directories.
"""

import argparse
import os
import sys
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed

try:
    from docx import Document
    from pptx import Presentation
except ImportError as e:
    print(f"Missing required package: {e}")
    print("\nPlease install required packages:")
    print("pip install --user python-docx python-pptx")
    sys.exit(1)


def docx_to_markdown(docx_path: Path) -> str:
    """Convert a DOCX file to Markdown."""
    try:
        doc = Document(str(docx_path))
        md_lines = []

        # Add title from filename
        md_lines.append(f"# {docx_path.stem}\n")
        md_lines.append("")

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                md_lines.append("")
                continue

            # Detect headings by style
            style_name = para.style.name.lower() if para.style else ""

            if 'heading 1' in style_name or 'title' in style_name:
                md_lines.append(f"## {text}\n")
            elif 'heading 2' in style_name:
                md_lines.append(f"### {text}\n")
            elif 'heading 3' in style_name:
                md_lines.append(f"#### {text}\n")
            elif 'list' in style_name or text.startswith(('•', '-', '●', '○', '■')):
                # Clean up bullet characters
                clean_text = text.lstrip('•-●○■◦▪ ')
                md_lines.append(f"- {clean_text}")
            else:
                md_lines.append(text)

        # Process tables
        for table in doc.tables:
            md_lines.append("")
            for i, row in enumerate(table.rows):
                cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                md_lines.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    md_lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
            md_lines.append("")

        return "\n".join(md_lines)

    except Exception as e:
        return f"# {docx_path.stem}\n\nError converting file: {e}\n"


def pptx_to_markdown(pptx_path: Path) -> str:
    """Convert a PPTX file to Markdown."""
    try:
        prs = Presentation(str(pptx_path))
        md_lines = []

        # Add title from filename
        md_lines.append(f"# {pptx_path.stem}\n")
        md_lines.append("")

        for slide_num, slide in enumerate(prs.slides, 1):
            md_lines.append(f"---\n")
            md_lines.append(f"## Slide {slide_num}\n")

            # Get title if present
            if slide.shapes.title and slide.shapes.title.text.strip():
                md_lines.append(f"### {slide.shapes.title.text.strip()}\n")

            # Extract text from shapes
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue  # Already handled

                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    for line in text.split('\n'):
                        line = line.strip()
                        if line:
                            # Check if it looks like a bullet point
                            if line.startswith(('•', '-', '●', '○', '■', '►', '▪')):
                                clean_line = line.lstrip('•-●○■►▪◦ ')
                                md_lines.append(f"- {clean_line}")
                            else:
                                md_lines.append(f"- {line}")

                # Handle tables
                if shape.has_table:
                    table = shape.table
                    md_lines.append("")
                    for i, row in enumerate(table.rows):
                        cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                        md_lines.append("| " + " | ".join(cells) + " |")
                        if i == 0:
                            md_lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
                    md_lines.append("")

            # Add speaker notes if present
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    md_lines.append(f"\n**Speaker Notes:** {notes_text}\n")

            md_lines.append("")

        return "\n".join(md_lines)

    except Exception as e:
        return f"# {pptx_path.stem}\n\nError converting file: {e}\n"


def convert_file(file_path: Path, output_dir: Path, format: str = "md") -> dict:
    """Convert a single file and return status."""
    result = {
        'file': file_path.name,
        'success': False,
        'output': None,
        'error': None
    }

    ext = file_path.suffix.lower()

    if ext == '.docx':
        content = docx_to_markdown(file_path)
    elif ext == '.pptx':
        content = pptx_to_markdown(file_path)
    elif ext == '.pdf':
        # PDFs are already readable by Claude Code
        result['success'] = True
        result['output'] = file_path
        result['skipped'] = True
        return result
    else:
        result['error'] = f"Unsupported format: {ext}"
        return result

    # Write output
    try:
        output_file = output_dir / f"{file_path.stem}.md"
        output_file.write_text(content, encoding='utf-8')
        result['success'] = True
        result['output'] = output_file
    except Exception as e:
        result['error'] = str(e)

    return result


def convert_directory(input_dir: Path, output_dir: Path = None,
                      format: str = "md", parallel: bool = True) -> dict:
    """Convert all supported files in a directory."""

    if output_dir is None:
        output_dir = input_dir / "generated"

    output_dir.mkdir(parents=True, exist_ok=True)

    # Find all convertible files
    supported_extensions = {'.docx', '.pptx'}
    files_to_convert = []

    for file_path in input_dir.iterdir():
        if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
            files_to_convert.append(file_path)

    print(f"\nFound {len(files_to_convert)} files to convert")
    print(f"Output directory: {output_dir}\n")

    stats = {
        'total': len(files_to_convert),
        'success': 0,
        'failed': 0,
        'skipped': 0,
        'errors': []
    }

    if parallel and len(files_to_convert) > 1:
        # Parallel conversion
        with ThreadPoolExecutor(max_workers=4) as executor:
            futures = {
                executor.submit(convert_file, f, output_dir, format): f
                for f in files_to_convert
            }

            for i, future in enumerate(as_completed(futures), 1):
                file_path = futures[future]
                result = future.result()

                if result.get('skipped'):
                    stats['skipped'] += 1
                    print(f"[{i}/{stats['total']}] Skipped: {result['file']} (already readable)")
                elif result['success']:
                    stats['success'] += 1
                    print(f"[{i}/{stats['total']}] Converted: {result['file']}")
                else:
                    stats['failed'] += 1
                    stats['errors'].append({'file': result['file'], 'error': result['error']})
                    print(f"[{i}/{stats['total']}] Failed: {result['file']} - {result['error']}")
    else:
        # Sequential conversion
        for i, file_path in enumerate(files_to_convert, 1):
            result = convert_file(file_path, output_dir, format)

            if result.get('skipped'):
                stats['skipped'] += 1
                print(f"[{i}/{stats['total']}] Skipped: {result['file']} (already readable)")
            elif result['success']:
                stats['success'] += 1
                print(f"[{i}/{stats['total']}] Converted: {result['file']}")
            else:
                stats['failed'] += 1
                stats['errors'].append({'file': result['file'], 'error': result['error']})
                print(f"[{i}/{stats['total']}] Failed: {result['file']} - {result['error']}")

    return stats


def main():
    parser = argparse.ArgumentParser(
        description='Convert DOCX and PPTX files to Markdown for Claude Code',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
Examples:
  # Convert all files in teks/practices to markdown
  python convert_to_readable.py teks/practices

  # Convert to a specific output directory
  python convert_to_readable.py teks/practices --output teks/practices_md

  # Convert sequentially (slower but uses less memory)
  python convert_to_readable.py teks/practices --no-parallel

Notes:
  - DOCX files are converted to Markdown with headings, lists, and tables
  - PPTX files are converted to Markdown with slide structure preserved
  - PDF files are skipped (Claude Code can already read them)
  - Output files are saved with .md extension
        '''
    )

    parser.add_argument(
        'input_dir',
        type=Path,
        help='Directory containing files to convert'
    )

    parser.add_argument(
        '--output', '-o',
        type=Path,
        default=None,
        help='Output directory (default: <input_dir>/generated)'
    )

    parser.add_argument(
        '--no-parallel',
        action='store_true',
        help='Disable parallel processing'
    )

    args = parser.parse_args()

    if not args.input_dir.exists():
        print(f"Error: Directory not found: {args.input_dir}")
        sys.exit(1)

    if not args.input_dir.is_dir():
        print(f"Error: Not a directory: {args.input_dir}")
        sys.exit(1)

    print("="*60)
    print("Document Converter")
    print("="*60)
    print(f"Input directory: {args.input_dir}")

    stats = convert_directory(
        args.input_dir,
        args.output,
        parallel=not args.no_parallel
    )

    print("\n" + "="*60)
    print("Conversion Summary")
    print("="*60)
    print(f"Total files:    {stats['total']}")
    print(f"Converted:      {stats['success']}")
    print(f"Skipped:        {stats['skipped']}")
    print(f"Failed:         {stats['failed']}")

    if stats['errors']:
        print("\nFailed files:")
        for err in stats['errors']:
            print(f"  - {err['file']}: {err['error']}")

    print("="*60)


if __name__ == "__main__":
    main()
