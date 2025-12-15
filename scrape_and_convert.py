#!/usr/bin/env python3
"""
Spider and scrape PowerPoint presentations from growtexasteachers.org
Convert to Markdown and PDF formats with various combination options.
"""

import argparse
import os
import re
import sys
from pathlib import Path
from urllib.parse import urljoin, urlparse
import time

try:
    import requests
    from bs4 import BeautifulSoup
    from pptx import Presentation
    import markdown
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT, TA_CENTER
except ImportError as e:
    print(f"Missing required package: {e}")
    print("\nPlease install required packages:")
    print("pip install --user requests beautifulsoup4 python-pptx markdown reportlab")
    sys.exit(1)


class PowerPointScraper:
    """Scrapes and converts PowerPoint presentations from growtexasteachers.org"""

    def __init__(self, base_dir="./output", delay=0.5):
        self.base_dir = Path(base_dir)
        self.delay = delay
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        })

    def scrape_page(self, url):
        """Scrape a page and extract PowerPoint links"""
        print(f"\nScraping: {url}")
        try:
            response = self.session.get(url, timeout=30)
            response.raise_for_status()
            soup = BeautifulSoup(response.content, 'html.parser')

            # Find all links that end with .pptx
            pptx_links = []
            for link in soup.find_all('a', href=True):
                href = link['href']
                if href.endswith('.pptx') or '.pptx?' in href:
                    full_url = urljoin(url, href)
                    # Extract a clean filename
                    filename = self._extract_filename(href)
                    pptx_links.append({
                        'url': full_url,
                        'filename': filename,
                        'text': link.get_text(strip=True)
                    })

            print(f"Found {len(pptx_links)} PowerPoint files")
            return pptx_links

        except Exception as e:
            print(f"Error scraping {url}: {e}")
            return []

    def _extract_filename(self, href):
        """Extract a clean filename from a URL"""
        # Remove query parameters
        href = href.split('?')[0]
        # Get the last part of the path
        filename = href.split('/')[-1]
        # Clean up the filename
        filename = re.sub(r'[^\w\-_\.]', '_', filename)
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        return filename

    def download_pptx(self, url, output_path):
        """Download a PowerPoint file"""
        try:
            print(f"Downloading: {output_path.name}")
            response = self.session.get(url, timeout=60)
            response.raise_for_status()

            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_bytes(response.content)

            time.sleep(self.delay)  # Be polite to the server
            return True

        except Exception as e:
            print(f"Error downloading {url}: {e}")
            return False

    def pptx_to_markdown(self, pptx_path):
        """Convert PowerPoint to Markdown"""
        try:
            prs = Presentation(str(pptx_path))
            markdown_lines = []

            # Add title
            markdown_lines.append(f"# {pptx_path.stem}\n")

            for slide_num, slide in enumerate(prs.slides, 1):
                markdown_lines.append(f"\n## Slide {slide_num}\n")

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        # Check if it might be a title (usually larger/bold)
                        if shape == slide.shapes.title:
                            markdown_lines.append(f"### {text}\n")
                        else:
                            # Format as bullet points or paragraphs
                            for line in text.split('\n'):
                                if line.strip():
                                    markdown_lines.append(f"- {line.strip()}\n")

                    # Extract tables
                    if shape.has_table:
                        table = shape.table
                        # Create markdown table
                        markdown_lines.append("\n")
                        for row_idx, row in enumerate(table.rows):
                            cells = [cell.text.strip() for cell in row.cells]
                            markdown_lines.append("| " + " | ".join(cells) + " |\n")
                            if row_idx == 0:
                                markdown_lines.append("| " + " | ".join(["---"] * len(cells)) + " |\n")
                        markdown_lines.append("\n")

                # Add notes if present
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        markdown_lines.append(f"\n**Notes:** {notes_text}\n")

            return "".join(markdown_lines)

        except Exception as e:
            print(f"Error converting {pptx_path} to markdown: {e}")
            return f"# {pptx_path.stem}\n\nError converting file: {e}\n"

    def markdown_to_pdf(self, markdown_text, output_path):
        """Convert Markdown to PDF using ReportLab"""
        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)

            # Create PDF document
            doc = SimpleDocTemplate(
                str(output_path),
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=18
            )

            # Container for the 'Flowable' objects
            elements = []

            # Define styles
            styles = getSampleStyleSheet()

            # Custom styles
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor=colors.HexColor('#2c3e50'),
                spaceAfter=30,
                spaceBefore=12
            )

            heading1_style = ParagraphStyle(
                'CustomHeading1',
                parent=styles['Heading1'],
                fontSize=18,
                textColor=colors.HexColor('#34495e'),
                spaceAfter=12,
                spaceBefore=12
            )

            heading2_style = ParagraphStyle(
                'CustomHeading2',
                parent=styles['Heading2'],
                fontSize=14,
                textColor=colors.HexColor('#7f8c8d'),
                spaceAfter=10,
                spaceBefore=10
            )

            body_style = ParagraphStyle(
                'CustomBody',
                parent=styles['BodyText'],
                fontSize=11,
                spaceAfter=12
            )

            bullet_style = ParagraphStyle(
                'CustomBullet',
                parent=styles['BodyText'],
                fontSize=11,
                leftIndent=20,
                spaceAfter=6
            )

            # Parse markdown and convert to PDF elements
            lines = markdown_text.split('\n')
            i = 0
            while i < len(lines):
                line = lines[i].strip()

                if not line:
                    i += 1
                    continue

                # Handle headers
                if line.startswith('# '):
                    elements.append(Paragraph(line[2:], title_style))
                elif line.startswith('## '):
                    elements.append(Paragraph(line[3:], heading1_style))
                elif line.startswith('### '):
                    elements.append(Paragraph(line[4:], heading2_style))

                # Handle bullet points
                elif line.startswith('- '):
                    elements.append(Paragraph(f"• {line[2:]}", bullet_style))

                # Handle horizontal rules
                elif line.startswith('---'):
                    elements.append(Spacer(1, 0.3*inch))
                    elements.append(PageBreak())

                # Handle markdown tables
                elif '|' in line and i + 1 < len(lines) and '---' in lines[i + 1]:
                    # Parse table
                    table_data = []
                    # Header row
                    header = [cell.strip() for cell in line.split('|')[1:-1]]
                    table_data.append(header)

                    # Skip separator line
                    i += 2

                    # Data rows
                    while i < len(lines) and '|' in lines[i] and '---' not in lines[i]:
                        row = [cell.strip() for cell in lines[i].split('|')[1:-1]]
                        if row:
                            table_data.append(row)
                        i += 1

                    # Create table
                    if table_data:
                        t = Table(table_data)
                        t.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#ecf0f1')),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 11),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.white),
                            ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                            ('FONTSIZE', (0, 1), (-1, -1), 10),
                            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#bdc3c7')),
                            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                        ]))
                        elements.append(t)
                        elements.append(Spacer(1, 0.2*inch))
                    i -= 1

                # Handle bold text (Notes, etc.)
                elif '**' in line:
                    # Simple bold text handling
                    text = line.replace('**', '<b>', 1).replace('**', '</b>', 1)
                    elements.append(Paragraph(text, body_style))

                # Regular paragraph
                else:
                    if line:
                        elements.append(Paragraph(line, body_style))

                i += 1

            # Build PDF
            doc.build(elements)
            return True

        except Exception as e:
            print(f"Error converting to PDF: {e}")
            import traceback
            traceback.print_exc()
            return False

    def process_url(self, url, subject_name):
        """Process a single URL: scrape, download, and convert"""
        print(f"\n{'='*60}")
        print(f"Processing: {subject_name}")
        print(f"{'='*60}")

        # Create directories
        subject_dir = self.base_dir / subject_name
        pptx_dir = subject_dir / "pptx"
        md_dir = subject_dir / "markdown"
        pdf_dir = subject_dir / "pdf"

        for d in [pptx_dir, md_dir, pdf_dir]:
            d.mkdir(parents=True, exist_ok=True)

        # Scrape for PowerPoint links
        pptx_links = self.scrape_page(url)

        if not pptx_links:
            print(f"No PowerPoint files found at {url}")
            return []

        # Download and convert each PowerPoint
        all_markdown = []
        successful_files = []

        for idx, link_info in enumerate(pptx_links, 1):
            print(f"\n[{idx}/{len(pptx_links)}] Processing: {link_info['filename']}")

            pptx_path = pptx_dir / link_info['filename']

            # Download
            if not pptx_path.exists():
                if not self.download_pptx(link_info['url'], pptx_path):
                    continue
            else:
                print(f"Already downloaded: {pptx_path.name}")

            # Convert to Markdown
            markdown_content = self.pptx_to_markdown(pptx_path)
            md_path = md_dir / f"{pptx_path.stem}.md"
            md_path.write_text(markdown_content, encoding='utf-8')
            print(f"Created: {md_path.name}")

            # Convert to PDF
            pdf_path = pdf_dir / f"{pptx_path.stem}.pdf"
            if self.markdown_to_pdf(markdown_content, pdf_path):
                print(f"Created: {pdf_path.name}")

            # Store for combined file
            all_markdown.append(markdown_content)
            all_markdown.append("\n\n---\n\n")  # Add separator
            successful_files.append(link_info['filename'])

        # Create combined files for this subject
        if all_markdown:
            combined_md = "".join(all_markdown)
            combined_md_path = subject_dir / f"{subject_name}_combined.md"
            combined_md_path.write_text(combined_md, encoding='utf-8')
            print(f"\nCreated combined markdown: {combined_md_path}")

            combined_pdf_path = subject_dir / f"{subject_name}_combined.pdf"
            if self.markdown_to_pdf(combined_md, combined_pdf_path):
                print(f"Created combined PDF: {combined_pdf_path}")

        return all_markdown

    def process_all(self, urls):
        """Process all URLs and create final combined files"""
        all_content = []

        for url, subject_name in urls:
            content = self.process_url(url, subject_name)
            all_content.extend(content)

        # Create final combined files in root directory
        if all_content:
            print(f"\n{'='*60}")
            print("Creating final combined files...")
            print(f"{'='*60}")

            final_md = "".join(all_content)

            # Combined markdown
            combined_md_path = self.base_dir / "all_classes_combined.md"
            combined_md_path.write_text(final_md, encoding='utf-8')
            print(f"Created: {combined_md_path}")

            # Combined PDF (all classes)
            combined_pdf_path = self.base_dir / "all_classes_combined.pdf"
            if self.markdown_to_pdf(final_md, combined_pdf_path):
                print(f"Created: {combined_pdf_path}")

            # Create combined PDF in root directory
            ct_pdf_path = Path.cwd() / "all_courses_combined.pdf"
            if self.markdown_to_pdf(final_md, ct_pdf_path):
                print(f"Created: {ct_pdf_path}")


def main():
    parser = argparse.ArgumentParser(
        description='Spider and convert PowerPoint files from growtexasteachers.org',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
Examples:
  # Process both default URLs
  python scrape_and_convert.py

  # Process with custom output directory
  python scrape_and_convert.py --output ./my_output

  # Process specific URLs
  python scrape_and_convert.py --url https://www.growtexasteachers.org/practices practices

  # Add delay between requests (be polite to server)
  python scrape_and_convert.py --delay 1.0
        '''
    )

    parser.add_argument(
        '--output', '-o',
        default='./output',
        help='Output directory for generated files (default: ./output)'
    )

    parser.add_argument(
        '--delay', '-d',
        type=float,
        default=0.5,
        help='Delay in seconds between requests (default: 0.5)'
    )

    parser.add_argument(
        '--url',
        nargs=2,
        action='append',
        metavar=('URL', 'NAME'),
        help='Add a URL to process with a subject name (can be used multiple times)'
    )

    args = parser.parse_args()

    # Default URLs if none specified
    if not args.url:
        urls = [
            ('https://www.growtexasteachers.org/practices', 'practices'),
            ('https://www.growtexasteachers.org/practicum', 'practicum')
        ]
    else:
        urls = args.url

    print("="*60)
    print("PowerPoint Scraper and Converter")
    print("="*60)
    print(f"Output directory: {args.output}")
    print(f"URLs to process: {len(urls)}")
    for url, name in urls:
        print(f"  - {name}: {url}")
    print("="*60)

    scraper = PowerPointScraper(base_dir=args.output, delay=args.delay)
    scraper.process_all(urls)

    print("\n" + "="*60)
    print("Processing complete!")
    print("="*60)


if __name__ == "__main__":
    main()
