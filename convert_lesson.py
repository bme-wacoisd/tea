"""
Convert lesson materials to PDF and PowerPoint formats.
"""
import re
from pathlib import Path
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem, Table, TableStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT, TA_CENTER
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def md_to_pdf(md_path: Path, pdf_path: Path):
    """Convert markdown lesson plan to PDF."""
    content = md_path.read_text(encoding='utf-8')

    doc = SimpleDocTemplate(
        str(pdf_path),
        pagesize=letter,
        rightMargin=0.75*inch,
        leftMargin=0.75*inch,
        topMargin=0.75*inch,
        bottomMargin=0.75*inch
    )

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        name='Title1',
        parent=styles['Heading1'],
        fontSize=18,
        spaceAfter=12,
        textColor=colors.HexColor('#1a365d')
    ))
    styles.add(ParagraphStyle(
        name='Heading2Custom',
        parent=styles['Heading2'],
        fontSize=14,
        spaceBefore=16,
        spaceAfter=8,
        textColor=colors.HexColor('#2c5282')
    ))
    styles.add(ParagraphStyle(
        name='Heading3Custom',
        parent=styles['Heading3'],
        fontSize=12,
        spaceBefore=12,
        spaceAfter=6,
        textColor=colors.HexColor('#2d3748')
    ))
    styles.add(ParagraphStyle(
        name='BodyCustom',
        parent=styles['Normal'],
        fontSize=10,
        spaceAfter=6,
        leading=14
    ))
    styles.add(ParagraphStyle(
        name='BulletCustom',
        parent=styles['Normal'],
        fontSize=10,
        leftIndent=20,
        spaceAfter=4,
        leading=14
    ))

    story = []
    lines = content.split('\n')
    i = 0

    while i < len(lines):
        line = lines[i].strip()

        # Skip horizontal rules
        if line == '---':
            story.append(Spacer(1, 12))
            i += 1
            continue

        # Title (h1)
        if line.startswith('# '):
            text = line[2:].strip()
            story.append(Paragraph(text, styles['Title1']))
            i += 1
            continue

        # H2
        if line.startswith('## '):
            text = line[3:].strip()
            story.append(Paragraph(text, styles['Heading2Custom']))
            i += 1
            continue

        # H3
        if line.startswith('### '):
            text = line[4:].strip()
            story.append(Paragraph(text, styles['Heading3Custom']))
            i += 1
            continue

        # Bullet points
        if line.startswith('- '):
            text = process_inline_formatting(line[2:].strip())
            story.append(Paragraph(f"• {text}", styles['BulletCustom']))
            i += 1
            continue

        # Numbered items
        num_match = re.match(r'^(\d+)\.\s+(.+)$', line)
        if num_match:
            num, text = num_match.groups()
            text = process_inline_formatting(text)
            story.append(Paragraph(f"{num}. {text}", styles['BulletCustom']))
            i += 1
            continue

        # Regular paragraph
        if line:
            text = process_inline_formatting(line)
            story.append(Paragraph(text, styles['BodyCustom']))

        i += 1

    doc.build(story)
    print(f"Created PDF: {pdf_path}")


def process_inline_formatting(text: str) -> str:
    """Convert markdown inline formatting to reportlab."""
    # Bold
    text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
    # Italic
    text = re.sub(r'\*(.+?)\*', r'<i>\1</i>', text)
    # Checkbox items
    text = text.replace('[ ]', '☐')
    text = text.replace('[x]', '☑')
    return text


def md_to_pptx(md_path: Path, pptx_path: Path, assets_path: Path = None):
    """Convert markdown slides to PowerPoint."""
    content = md_path.read_text(encoding='utf-8')

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Split by slide delimiter
    slides_content = content.split('\n---\n')

    for slide_idx, slide_content in enumerate(slides_content):
        slide_content = slide_content.strip()
        if not slide_content:
            continue

        lines = slide_content.split('\n')

        # Extract title
        title = ""
        body_lines = []
        speaker_notes = ""
        in_table = False
        table_rows = []

        for line in lines:
            line_stripped = line.strip()

            # Skip image references for now
            if line_stripped.startswith('!['):
                continue

            if line_stripped.startswith('Speaker notes:'):
                speaker_notes = line_stripped.replace('Speaker notes:', '').strip()
                continue

            if line_stripped.startswith('# '):
                title = line_stripped[2:].strip()
                continue

            if line_stripped.startswith('## '):
                title = line_stripped[3:].strip()
                continue

            # Handle tables
            if '|' in line_stripped and not line_stripped.startswith('|---'):
                if '---' in line_stripped:
                    continue  # Skip table header separator
                cells = [c.strip() for c in line_stripped.split('|') if c.strip()]
                if cells:
                    table_rows.append(cells)
                in_table = True
                continue
            elif in_table and not line_stripped:
                in_table = False

            if line_stripped and not in_table:
                body_lines.append(line_stripped)

        # Create slide
        if slide_idx == 0:
            # Title slide
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)

            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(26, 54, 93)
            p.alignment = PP_ALIGN.CENTER

            # Add subtitle if present
            if body_lines:
                subtitle = ' '.join(body_lines[:2])
                sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
                tf = sub_box.text_frame
                p = tf.paragraphs[0]
                p.text = clean_markdown(subtitle)
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor(74, 85, 104)
                p.alignment = PP_ALIGN.CENTER
        else:
            # Content slide
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)

            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(26, 54, 93)

            # Add body content
            content_top = Inches(1.4)

            if table_rows:
                # Create table
                table_width = Inches(11)
                row_height = Inches(0.4)
                num_rows = len(table_rows)
                num_cols = len(table_rows[0]) if table_rows else 2

                table = slide.shapes.add_table(
                    num_rows, num_cols,
                    Inches(1), content_top,
                    table_width, row_height * num_rows
                ).table

                # Style table
                for row_idx, row_data in enumerate(table_rows):
                    for col_idx, cell_text in enumerate(row_data):
                        if col_idx < num_cols:
                            cell = table.cell(row_idx, col_idx)
                            cell.text = clean_markdown(cell_text)
                            para = cell.text_frame.paragraphs[0]
                            para.font.size = Pt(14)
                            if row_idx == 0:
                                para.font.bold = True
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(44, 82, 130)
                                para.font.color.rgb = RGBColor(255, 255, 255)

            # Add bullet points
            if body_lines:
                body_box = slide.shapes.add_textbox(
                    Inches(0.5),
                    content_top + (Inches(2.5) if table_rows else Inches(0)),
                    Inches(12.333),
                    Inches(5)
                )
                tf = body_box.text_frame
                tf.word_wrap = True

                first_para = True
                for line in body_lines:
                    if first_para:
                        p = tf.paragraphs[0]
                        first_para = False
                    else:
                        p = tf.add_paragraph()

                    text = clean_markdown(line)

                    # Handle bullet points
                    if line.startswith('- ') or line.startswith('* '):
                        p.text = "• " + text[2:]
                        p.level = 0
                    elif line.startswith('  - ') or line.startswith('  * '):
                        p.text = "  ◦ " + text[4:]
                        p.level = 1
                    elif re.match(r'^\d+\.\s', line):
                        p.text = text
                        p.level = 0
                    elif line.startswith('> '):
                        p.text = '"' + text[2:] + '"'
                        p.font.italic = True
                    elif line.startswith('❌') or line.startswith('✅'):
                        p.text = text
                    else:
                        p.text = text

                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(45, 55, 72)
                    p.space_after = Pt(8)

        # Add speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    prs.save(str(pptx_path))
    print(f"Created PowerPoint: {pptx_path}")


def clean_markdown(text: str) -> str:
    """Remove markdown formatting from text."""
    # Remove bold
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    # Remove italic
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    # Remove links
    text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)
    # Clean up bullet markers if at start
    if text.startswith('- '):
        text = text[2:]
    return text.strip()


def main():
    lesson_dir = Path("lessons/01-james-baldwin-civil-rights")

    # Convert lesson plan to PDF
    lesson_plan_md = lesson_dir / "lesson-plan.md"
    lesson_plan_pdf = lesson_dir / "lesson-plan.pdf"
    md_to_pdf(lesson_plan_md, lesson_plan_pdf)

    # Convert slides to PowerPoint
    slides_md = lesson_dir / "slides.md"
    slides_pptx = lesson_dir / "slides.pptx"
    assets_dir = lesson_dir / "assets"
    md_to_pptx(slides_md, slides_pptx, assets_dir)


if __name__ == "__main__":
    main()
